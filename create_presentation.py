from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Professional academic theme
NAVY_BLUE = RgbColor(0, 51, 102)
DARK_BLUE = RgbColor(0, 32, 64)
WHITE = RgbColor(255, 255, 255)
LIGHT_GRAY = RgbColor(240, 240, 240)
BLACK = RgbColor(0, 0, 0)

def add_title_shape(slide, title_text, subtitle_text=None):
    """Add a styled title to the slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title_text
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY_BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Add underline bar  
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY_BLUE
    shape.line.fill.background()

def add_bullet_point(text_frame, text, level=0, font_size=18):
    """Add a bullet point to text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(font_size)
    p.font.color.rgb = BLACK
    p.space_before = Pt(6)
    p.space_after = Pt(3)

# ==================== SLIDE 1: Title Slide ====================
slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Background (light gradient effect using rectangle)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = WHITE
bg.line.fill.background()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.5))
tf = title_box.text_frame
tf.paragraphs[0].text = "Project Management and Team Collaboration Tool"
tf.paragraphs[0].font.size = Pt(40)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = NAVY_BLUE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Technology Stack
tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(0.6))
tf = tech_box.text_frame
tf.paragraphs[0].text = "Technology Stack: React.js | Java Spring Boot | MySQL"
tf.paragraphs[0].font.size = Pt(22)
tf.paragraphs[0].font.color.rgb = DARK_BLUE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Decorative line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4), Inches(5.333), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = NAVY_BLUE
shape.line.fill.background()

# Team Members
team_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(6), Inches(1.5))
tf = team_box.text_frame
tf.paragraphs[0].text = "Team Members:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
p = tf.add_paragraph()
p.text = "• [Student Name 1] - [Roll Number]"
p.font.size = Pt(16)
p.font.color.rgb = BLACK
p = tf.add_paragraph()
p.text = "• [Student Name 2] - [Roll Number]"
p.font.size = Pt(16)
p.font.color.rgb = BLACK
p = tf.add_paragraph()
p.text = "• [Student Name 3] - [Roll Number]"
p.font.size = Pt(16)
p.font.color.rgb = BLACK

# Guide Name
guide_box = slide.shapes.add_textbox(Inches(7), Inches(4.5), Inches(5.833), Inches(1))
tf = guide_box.text_frame
tf.paragraphs[0].text = "Project Guide:"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = BLACK
p = tf.add_paragraph()
p.text = "[Guide Name]"
p.font.size = Pt(16)
p.font.color.rgb = BLACK
p = tf.add_paragraph()
p.text = "[Designation]"
p.font.size = Pt(14)
p.font.color.rgb = DARK_BLUE

# College/Department
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
tf = footer_box.text_frame
tf.paragraphs[0].text = "[College Name] | [Department] | [Academic Year]"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = DARK_BLUE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ==================== SLIDE 2: Abstract ====================
slide = prs.slides.add_slide(slide_layout)
add_title_shape(slide, "Abstract")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

# Problem Overview
p = tf.paragraphs[0]
p.text = "Problem Overview"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE
p.space_after = Pt(8)

add_bullet_point(tf, "Organizations face challenges in managing multiple projects, tracking task progress, and ensuring effective team collaboration across distributed teams.")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Purpose of the Project"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE
p.space_after = Pt(8)

add_bullet_point(tf, "To develop a comprehensive web-based project management system that streamlines task allocation, progress monitoring, and team communication.")

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Solution Approach"
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE
p.space_after = Pt(8)

add_bullet_point(tf, "A full-stack web application built with React.js frontend and Java Spring Boot backend, utilizing MySQL for persistent data storage.")
add_bullet_point(tf, "JWT-based authentication ensures secure access, while role-based authorization enables appropriate access control for users, managers, and administrators.")
add_bullet_point(tf, "RESTful API architecture enables seamless communication between frontend and backend components.")

# ==================== SLIDE 3: Objective ====================
slide = prs.slides.add_slide(slide_layout)
add_title_shape(slide, "Objectives")

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True
tf.paragraphs[0].text = ""

objectives = [
    ("Task Management", "To provide an intuitive interface for creating, assigning, and tracking tasks with priorities, deadlines, and status updates."),
    ("Team Collaboration", "To enable effective team collaboration through project-based member management and role assignments."),
    ("Project Organization", "To develop a centralized system for organizing multiple projects with status tracking and progress visualization."),
    ("User Authentication", "To implement secure user authentication using JWT tokens and role-based access control (USER, MANAGER, ADMIN)."),
    ("Efficiency and Usability", "To improve overall project workflow efficiency through a responsive, user-friendly web interface."),
    ("Administrative Control", "To provide administrative capabilities for user management, role assignment, and system-wide project oversight.")
]

for i, (title, desc) in enumerate(objectives):
    if i > 0:
        p = tf.add_paragraph()
        p.text = ""
        p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = f"{i+1}. {title}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY_BLUE
    
    p = tf.add_paragraph()
    p.text = f"    {desc}"
    p.font.size = Pt(16)
    p.font.color.rgb = BLACK
    p.space_before = Pt(3)

# ==================== SLIDE 4: Literature Survey ====================
slide = prs.slides.add_slide(slide_layout)
add_title_shape(slide, "Literature Survey")

# Left column - Existing Tools
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5))
tf = left_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Existing Project Management Tools"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE

tools = [
    ("Jira (Atlassian)", "Enterprise-grade, complex setup, high learning curve"),
    ("Trello", "Kanban-based, limited for complex workflows"),
    ("Asana", "Feature-rich but requires paid plans for advanced features"),
    ("Microsoft Project", "Desktop-focused, expensive licensing")
]

for tool, desc in tools:
    p = tf.add_paragraph()
    p.text = f"• {tool}"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.space_before = Pt(8)
    
    p = tf.add_paragraph()
    p.text = f"   {desc}"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BLUE

# Right column - Technologies & Limitations
right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(5.5))
tf = right_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Common Technologies Used"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE

techs = ["Frontend: React, Angular, Vue.js", "Backend: Node.js, Java, Python", "Databases: PostgreSQL, MongoDB, MySQL", "Authentication: OAuth, JWT, SAML"]
for tech in techs:
    p = tf.add_paragraph()
    p.text = f"• {tech}"
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(16)

p = tf.add_paragraph()
p.text = "Identified Limitations"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE

limitations = ["High cost for small teams", "Complex configuration requirements", "Limited customization options", "Steep learning curves"]
for lim in limitations:
    p = tf.add_paragraph()
    p.text = f"• {lim}"
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_before = Pt(4)

# ==================== SLIDE 5: Existing System ====================
slide = prs.slides.add_slide(slide_layout)
add_title_shape(slide, "Existing System")

# Left Panel - Description
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5))
tf = left_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Traditional Project Management Methods"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE

methods = [
    "Manual tracking using spreadsheets (Excel/Google Sheets)",
    "Email-based task assignment and status updates",
    "Physical documentation and progress reports",
    "Standalone desktop applications without collaboration features",
    "Informal verbal communication for task delegation"
]
for method in methods:
    p = tf.add_paragraph()
    p.text = f"• {method}"
    p.font.size = Pt(15)
    p.font.color.rgb = BLACK
    p.space_before = Pt(8)

# Right Panel - Problems
right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(5.5))
tf = right_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Problems Identified"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE

problems = [
    ("Lack of Real-time Tracking", "No centralized visibility into task progress"),
    ("Poor Collaboration", "Difficulty in coordinating distributed team members"),
    ("Data Inconsistency", "Multiple versions of documents causing confusion"),
    ("Limited Accessibility", "Information tied to specific devices/locations"),
    ("Inefficient Reporting", "Manual effort required for progress reports"),
    ("No Role-based Access", "Security and permission management challenges")
]
for title, desc in problems:
    p = tf.add_paragraph()
    p.text = f"• {title}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = f"   {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_BLUE

# ==================== SLIDE 6: Proposed System ====================
slide = prs.slides.add_slide(slide_layout)
add_title_shape(slide, "Proposed System")

# Left Panel - Description
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.2))
tf = left_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Web-Based Solution"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE

p = tf.add_paragraph()
p.text = "A modern, full-stack web application providing comprehensive project management and team collaboration capabilities."
p.font.size = Pt(15)
p.font.color.rgb = BLACK
p.space_before = Pt(8)

p = tf.add_paragraph()
p.text = ""
p.space_before = Pt(10)

p = tf.add_paragraph()
p.text = "Key Features"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE

features = [
    "User Authentication with JWT-based security",
    "Project Management with CRUD operations",
    "Task Tracking with priorities and deadlines",
    "Team Collaboration with role-based access",
    "Admin Dashboard for system management",
    "Progress Visualization with metrics"
]
for feature in features:
    p = tf.add_paragraph()
    p.text = f"✓ {feature}"
    p.font.size = Pt(14)
    p.font.color.rgb = BLACK
    p.space_before = Pt(5)

# Right Panel - Advantages
right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(5.2))
tf = right_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Advantages Over Existing Systems"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = NAVY_BLUE

advantages = [
    ("Centralized Data", "All project information in one accessible location"),
    ("Real-time Updates", "Instant visibility into task and project status"),
    ("Secure Access", "JWT authentication with role-based permissions"),
    ("Cross-platform", "Accessible from any device with a web browser"),
    ("Scalability", "Designed to handle growing teams and projects"),
    ("Cost-effective", "Open-source technologies reduce licensing costs")
]
for title, desc in advantages:
    p = tf.add_paragraph()
    p.text = f"• {title}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.space_before = Pt(6)
    
    p = tf.add_paragraph()
    p.text = f"   {desc}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_BLUE

# ==================== SLIDE 7: Block Diagram ====================
slide = prs.slides.add_slide(slide_layout)
add_title_shape(slide, "System Architecture - Block Diagram")

# Try to add the generated block diagram image
image_path = r"C:\Users\Ritikesh\.gemini\antigravity\brain\b4421934-838a-4533-b9a0-13f41ebab9e0\block_diagram_1767824485799.png"
if os.path.exists(image_path):
    slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5))
else:
    # Fallback: Create text-based diagram representation
    diagram_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(4))
    tf = diagram_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "[Block Diagram Image - Please insert the architecture diagram manually]"
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

# Architecture description
desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.8))
tf = desc_box.text_frame
tf.paragraphs[0].text = "Three-tier Architecture: React Frontend → Spring Boot REST API → MySQL Database"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].font.color.rgb = DARK_BLUE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save the presentation
output_path = r"c:\Users\Ritikesh\Downloads\Full Stack development_Project Management Using JAVA + React\Project_Management_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
